import os
import shutil
import tempfile
import subprocess
from concurrent.futures import ThreadPoolExecutor

from pptx import Presentation

from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdftypes import PDFObjectNotFound
from pdfminer import pdfinterp


class PDFConverter:
    def __init__(self, pdfname):
        self.presentation = Presentation()
        self._layout = self.presentation.slide_layouts[6]  # Blank slide
        self._width = self.presentation.slide_width
        self._height = self.presentation.slide_height
        self._pdfname = pdfname
        self._tempdir = tempfile.mkdtemp()

    def _task(self, i):
        cmd = """\
        convert \
        -density 600 \
        -colorspace sRGB \
        -background white \
        -alpha remove \
        -resize 1280x1024 \
        {0}[{2}] \
        {1}/page-{2}.png \
        """.format(self._pdfname, self._tempdir, i)

        subprocess.call(cmd, shell=True)

    def convert(self):
        with open(self._pdfname, 'rb') as f:
            parser = PDFParser(f)
            document = PDFDocument(parser, '')
            n_pages = pdfinterp.resolve1(document.catalog['Pages'])['Count']

        with ThreadPoolExecutor() as executor:
            executor.map(self._task, range(n_pages))

        for i in range(n_pages):
            slide = self.presentation.slides.add_slide(self._layout)
            slide.shapes.add_picture(
                '{}/page-{}.png'.format(self._tempdir, i),
                0,
                0,
                self._width,
                self._height,
            )

        shutil.rmtree(self._tempdir)

        return self.presentation


class PDFCommentExtracter:
    def __init__(self, pdfname):
        self._pdfname = pdfname
        self._pages = list()
        self.comments = dict()

    def _extract(self, objid, obj):
        if 'Type' in obj:
            if obj['Type'].name == 'Page':
                self._pages.append(objid)
            elif obj['Type'].name == 'Annot':
                if 'Subtype' in obj and 'Name' in obj \
                   and obj['Subtype'].name == 'Text' \
                   and obj['Name'].name == 'Comment':
                    page = len(self._pages) + 1
                    comment = obj['Contents'].decode('utf-8').strip()
                    self.comments[page] = comment

    def extract(self):
        with open(self._pdfname, 'rb') as f:
            parser = PDFParser(f)
            document = PDFDocument(parser, '')
            visited = set()

            for xref in document.xrefs:
                for objid in xref.get_objids():
                    if objid in visited:
                        continue
                    visited.add(objid)
                    try:
                        obj = document.getobj(objid)
                        if not isinstance(obj, dict):
                            continue
                        self._extract(objid, obj)
                    except PDFObjectNotFound as e:
                        pass

        return self.comments


class NoteInjector:
    def __init__(self, presentation, comments):
        self.presentation = presentation
        self._comments = comments

    def inject(self):
        for i, slide in enumerate(self.presentation.slides, 1):
            if i in self._comments:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = self._comments[i]

        return self.presentation


if __name__ == '__main__':
    import sys

    if len(sys.argv) > 1:
        pdfname = sys.argv[1]
    else:
        pdfname = 'slide.pdf'

    pptxname = pdfname + '.pptx'

    comments = PDFCommentExtracter(pdfname).extract()
    presentation = PDFConverter(pdfname).convert()
    presentation_with_notes = NoteInjector(presentation, comments).inject()

    presentation_with_notes.save(pptxname)
